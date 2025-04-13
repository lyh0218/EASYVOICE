import os
import time
from pptx import Presentation
from pptx.util import Inches
from pydub import AudioSegment

ffmpeg_path = "ffmpeg/"
os.environ["PATH"] = ffmpeg_path + os.pathsep + os.environ["PATH"]

def wavs_inference(tts_model,config):
    print(config["target_rms"] is float)
    fix_duration_value = float(config["fix_duration"])
    if fix_duration_value == 0:
        fix_duration_value = None
    new_ref_text = config["ref_text"]
    new_ref_text = f"。。{new_ref_text}。。"
    wav, sr, spec = tts_model.infer(
        ref_file = config["ref_file"],
        ref_text = new_ref_text,
        gen_text = config["gen_text"],
        file_wave = config["file_wave"],
        show_info = print,
        progress = None,
        target_rms = float(config["target_rms"]),
        speed = float(config["speed"]),
        cfg_strength = int(config["cfg_strength"]),
        nfe_step = int(config["nfe_step"]),
        sway_sampling_coef = int(config["sway_sampling_coef"]),
        cross_fade_duration = float(config["cross_fade_duration"]),
        fix_duration = fix_duration_value,
        remove_silence = config["remove_silence"],
        seed = None,
        file_spec = None
    )
    return config["file_wave"]

def ppt_fast_inference(tts_model,config,input_path,output_path):
    prs = Presentation(input_path)
    config["file_wave"] = "data/temp/ppt_voice_temp.wav"
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content = f"。。{shape.text_frame.text}。。"
                config["gen_text"] = text_content
                audio_path = wavs_inference(tts_model,config)
                left = shape.left
                top = shape.top
                width = Inches(0.5)
                height = Inches(0.5)
                audio_left = left - Inches(0.2)
                audio_top = top - Inches(0.2)
                slide.shapes.add_movie(audio_path, audio_left, audio_top, width, height, mime_type='audio/wav')
    prs.save(output_path)
    return output_path

def count_textboxes(ppt_path):
    prs = Presentation(ppt_path)
    text_num = 0
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_num += 1
    return text_num

def get_ppt_text(index,ppt_path):
    prs = Presentation(ppt_path)
    textbox_count = 0
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                textbox_count += 1
                text_content = shape.text_frame.text
                if textbox_count == index:
                   return text_content

def insert_audio_to_ppt(temp_json):
    prs = Presentation("data/temp/ppt_temp.pptx")
    target_num = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                target_num += 1
                if temp_json[f"shape_{target_num}"] != "NOT":
                    left = shape.left
                    top = shape.top
                    width = Inches(0.5)
                    height = Inches(0.5)
                    audio_left = left - Inches(0.2)
                    audio_top = top - Inches(0.2)
                    slide.shapes.add_movie(temp_json[f"shape_{target_num}"], audio_left, audio_top, width, height, mime_type='audio/wav')
    prs.save("data/temp/ppt_temp.pptx")
    return "data/temp/ppt_temp.pptx"